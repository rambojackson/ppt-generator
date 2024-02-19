from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_chart_slide(prs, title, chart_data, chart_type, legend_position):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    chart = slide.shapes.add_chart(chart_type, Inches(1), Inches(1.5), Inches(6), Inches(4.5), chart_data).chart
    chart.has_legend = True
    chart.legend.position = legend_position

def add_table_slide(prs, title, table_data):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    table = slide.shapes.add_table(rows=len(table_data), cols=len(table_data[0]), left=Inches(1), top=Inches(1.5), width=Inches(6), height=Inches(2)).table
    for i, row in enumerate(table_data):
        for j, cell in enumerate(row):
            table.cell(i, j).text = str(cell)
            table.cell(i, j).text_frame.paragraphs[0].font.size = Inches(0.15)

def add_image_slide(prs, title, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(6), height=Inches(4))

def add_text_box_slide(prs, title, text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(4))
    textbox.text_frame.text = text

def add_rectangle_slide(prs, title, color=RGBColor(255, 0, 0)):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(2), Inches(2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color

# Create a presentation object
prs = Presentation()

# Add title slide
add_title_slide(prs, "Sample Presentation", "Generated using Python \n - Rambo Jackson")

# Add chart slides
bar_chart_data = CategoryChartData()
bar_chart_data.categories = ['East', 'West', 'Midwest']
bar_chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
add_chart_slide(prs, "Bar Chart", bar_chart_data, XL_CHART_TYPE.COLUMN_CLUSTERED, XL_LEGEND_POSITION.TOP)

line_chart_data = CategoryChartData()
line_chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
line_chart_data.add_series('Series 1', (3.5, 2.7, 4.8, 5.2))
line_chart_data.add_series('Series 2', (2.6, 3.2, 2.4, 2.9))
add_chart_slide(prs, "Line Chart", line_chart_data, XL_CHART_TYPE.LINE, XL_LEGEND_POSITION.BOTTOM)

pie_chart_data = CategoryChartData()
pie_chart_data.categories = ['Apple', 'Banana', 'Orange']
pie_chart_data.add_series('Series 1', (40, 30, 30))
add_chart_slide(prs, "Pie Chart", pie_chart_data, XL_CHART_TYPE.PIE, XL_LEGEND_POSITION.RIGHT)

# Add table slide
table_data = [['ID', 'Name', 'Age'], [1, 'Alice', 30], [2, 'Bob', 35], [3, 'Charlie', 25]]
add_table_slide(prs, "Table", table_data)

# Add image slide
add_image_slide(prs, "Image", 'image.jpg')

# Add text box slide
add_text_box_slide(prs, "Text Box", "This is a text box. You can add text, bullet points, and more!")

# Add rectangle slide
add_rectangle_slide(prs, "Rectangle")

# Save the presentation
prs.save('sample_presentation.pptx')
